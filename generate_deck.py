#!/usr/bin/env python3
"""
ClosetFinder PM Portfolio Deck Generator
Generates a professional 14-slide PowerPoint deck for Julia Lee's PM portfolio.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color constants ──────────────────────────────────────────────────────────
BLUE        = RGBColor(0,   122, 204)   # #007ACC  primary
NAVY        = RGBColor(26,   26,  46)   # #1A1A2E  dark nav
WHITE       = RGBColor(255, 255, 255)
LIGHT_GRAY  = RGBColor(245, 245, 245)
GREEN       = RGBColor(46,  125,  50)
BODY_TEXT   = RGBColor(51,   51,  51)
LIGHT_BLUE  = RGBColor(210, 234, 255)   # box fill helper
DARK_BLUE   = RGBColor(0,    90, 160)   # darker shade helper

# ── Slide dimensions (widescreen 13.33 × 7.5 in) ────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Helper utilities ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    """Fill slide background with a solid color rectangle."""
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=BODY_TEXT,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect_with_text(slide, left, top, width, height,
                       fill_color, text, font_size=13,
                       text_color=WHITE, bold=False,
                       line_color=None, align=PP_ALIGN.CENTER,
                       v_anchor=None):
    from pptx.enum.text import MSO_ANCHOR
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    if v_anchor:
        tf.auto_size = None
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def add_bottom_border(slide):
    """Thin blue line at the bottom of non-blue slides."""
    bar = slide.shapes.add_shape(1, 0, H - Inches(0.07), W, Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def slide_title(slide, title_text, top=Inches(0.3), color=NAVY):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = color
    return txBox


def add_divider(slide, top, color=BLUE):
    bar = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_bullet_paragraph(tf, text, font_size=15, color=BODY_TEXT,
                         bold=False, space_before=None, indent=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.text = ""
    if space_before:
        p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if indent:
        p.level = 1
    return p


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BLUE)

    # Main title
    add_textbox(slide,
                "AI Stylist From Your Closet",
                Inches(1), Inches(2.0), Inches(11.33), Inches(1.4),
                font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide,
                "Your existing wardrobe, curated by AI",
                Inches(1), Inches(3.6), Inches(11.33), Inches(0.7),
                font_size=22, bold=False, color=WHITE, align=PP_ALIGN.CENTER,
                italic=True)

    # Thin white divider
    bar = slide.shapes.add_shape(1, Inches(3.5), Inches(4.5), Inches(6.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()

    # Bottom line
    add_textbox(slide,
                "Julia Lee  ·  Product Requirements Document  ·  March 2026  ·  CONFIDENTIAL",
                Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "The Problem")
    add_divider(slide, Inches(1.05))

    problems = [
        "Getting dressed is a daily decision tax — the average person spends 17 minutes choosing an outfit",
        "Wardrobes are full, yet people feel they have 'nothing to wear' — 80% of owned clothes go unworn",
        "Shopping is reactive and uncoordinated — new purchases rarely integrate with existing wardrobe",
    ]

    top_positions = [Inches(1.5), Inches(3.0), Inches(4.5)]

    for i, (prob, top) in enumerate(zip(problems, top_positions)):
        # Numbered circle
        circle = slide.shapes.add_shape(9,  # oval
            Inches(0.5), top + Inches(0.1), Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = WHITE

        add_textbox(slide, prob,
                    Inches(1.1), top, Inches(11.5), Inches(0.9),
                    font_size=16, color=BODY_TEXT)


def slide_03_solution(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "The Solution")
    add_divider(slide, Inches(1.05))

    # Hero statement
    add_textbox(slide,
                "An AI personal stylist that works with clothes you already own.",
                Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.7),
                font_size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Three value prop boxes
    boxes = [
        ("Digitize Your Closet",
         "Import your wardrobe automatically from retailer accounts or upload photos"),
        ("AI-Generated Outfits",
         "Context-aware looks built from your actual items — weather, occasion, calendar"),
        ("Discover What's Missing",
         "Fill closet gaps with targeted suggestions that genuinely match your style"),
    ]

    box_w = Inches(3.8)
    box_h = Inches(3.2)
    gap   = Inches(0.3)
    start_left = Inches(0.55)
    top   = Inches(2.1)

    for i, (title, body) in enumerate(boxes):
        left = start_left + i * (box_w + gap)

        # Box background
        bg = slide.shapes.add_shape(1, left, top, box_w, box_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = BLUE

        # Header strip
        hdr = slide.shapes.add_shape(1, left, top, box_w, Inches(0.55))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = BLUE
        hdr.line.fill.background()
        tf = hdr.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE

        add_textbox(slide, body,
                    left + Inches(0.15), top + Inches(0.65),
                    box_w - Inches(0.3), box_h - Inches(0.75),
                    font_size=14, color=BODY_TEXT, wrap=True)


def slide_04_how_it_works(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "How It Works")
    add_divider(slide, Inches(1.05))

    steps = [
        ("1", "Sign Up & Connect",
         "Create account, connect Google or Apple shopping account"),
        ("2", "Build Your Closet",
         "Select brands, connect retailer accounts, digitize wardrobe automatically"),
        ("3", "Get Styled by AI",
         "AI generates 3–6 outfit options using weather, day, and occasion context"),
        ("4", "Shop the Gaps",
         "Targeted product suggestions fill missing slots with affiliate-tracked links"),
    ]

    box_w = Inches(2.8)
    box_h = Inches(3.6)
    gap   = Inches(0.25)
    arrow_w = Inches(0.3)
    total = 4 * box_w + 3 * (gap + arrow_w)
    start_left = (W - total) / 2
    top = Inches(1.7)

    for i, (num, title, body) in enumerate(steps):
        left = start_left + i * (box_w + gap + arrow_w)

        # Main box
        bg = slide.shapes.add_shape(1, left, top, box_w, box_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = BLUE

        # Number circle
        circ = slide.shapes.add_shape(9, left + Inches(0.1), top + Inches(0.1),
                                      Inches(0.5), Inches(0.5))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        tf = circ.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Title
        add_textbox(slide, title,
                    left + Inches(0.1), top + Inches(0.65),
                    box_w - Inches(0.2), Inches(0.6),
                    font_size=14, bold=True, color=NAVY)

        # Body
        add_textbox(slide, body,
                    left + Inches(0.1), top + Inches(1.3),
                    box_w - Inches(0.2), Inches(2.0),
                    font_size=13, color=BODY_TEXT, wrap=True)

        # Arrow between steps
        if i < 3:
            arr_left = left + box_w + gap
            arr_top  = top + box_h / 2 - Inches(0.2)
            add_textbox(slide, "→",
                        arr_left, arr_top, arrow_w, Inches(0.4),
                        font_size=22, bold=True, color=BLUE,
                        align=PP_ALIGN.CENTER)


def slide_05_personas(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "Who We're Building For")
    add_divider(slide, Inches(1.05))

    personas = [
        {
            "tag": "Primary Consumer",
            "name": "The Style-Conscious Professional",
            "meta": "Age: 25–42",
            "pain_header": "Pain Points",
            "pain": (
                "• Knows what she likes but not what to wear each day.\n"
                "• Owns plenty but feels wardrobe is underused.\n"
                "• Impulse-buys items that don't coordinate."
            ),
            "goals_header": "Goals",
            "goals": (
                "• Look put-together daily with minimal effort.\n"
                "• Get more value from existing wardrobe."
            ),
        },
        {
            "tag": "B2B User",
            "name": "Fashion Brand Product Manager",
            "meta": "Role: AI Product Designer Platform Client",
            "pain_header": "Need",
            "pain": (
                "Real consumer purchase and preference data to inform product design"
            ),
            "goals_header": "Value",
            "goals": (
                "Validated taxonomy, trend signals, lifestyle cohort profiles from real behavior"
            ),
        },
    ]

    box_w = Inches(5.9)
    box_h = Inches(5.4)
    top   = Inches(1.4)
    lefts = [Inches(0.5), Inches(6.9)]

    for i, (p, left) in enumerate(zip(personas, lefts)):
        # Outer box
        bg = slide.shapes.add_shape(1, left, top, box_w, box_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = BLUE

        # Header strip
        hdr = slide.shapes.add_shape(1, left, top, box_w, Inches(0.5))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = BLUE
        hdr.line.fill.background()
        tf = hdr.text_frame
        ph = tf.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        run = ph.add_run()
        run.text = p["tag"]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Name
        add_textbox(slide, p["name"],
                    left + Inches(0.15), top + Inches(0.6),
                    box_w - Inches(0.3), Inches(0.55),
                    font_size=15, bold=True, color=NAVY)

        # Meta
        add_textbox(slide, p["meta"],
                    left + Inches(0.15), top + Inches(1.15),
                    box_w - Inches(0.3), Inches(0.4),
                    font_size=13, color=BODY_TEXT, italic=True)

        # Pain header
        add_textbox(slide, p["pain_header"],
                    left + Inches(0.15), top + Inches(1.65),
                    box_w - Inches(0.3), Inches(0.35),
                    font_size=13, bold=True, color=BLUE)

        # Pain text
        add_textbox(slide, p["pain"],
                    left + Inches(0.15), top + Inches(2.0),
                    box_w - Inches(0.3), Inches(1.4),
                    font_size=13, color=BODY_TEXT, wrap=True)

        # Goals header
        add_textbox(slide, p["goals_header"],
                    left + Inches(0.15), top + Inches(3.55),
                    box_w - Inches(0.3), Inches(0.35),
                    font_size=13, bold=True, color=BLUE)

        # Goals text
        add_textbox(slide, p["goals"],
                    left + Inches(0.15), top + Inches(3.9),
                    box_w - Inches(0.3), Inches(1.2),
                    font_size=13, color=BODY_TEXT, wrap=True)


def slide_06_architecture(prs):
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_bottom_border(slide)
    slide_title(slide, "Product Architecture")
    add_divider(slide, Inches(1.05))

    layers = [
        (BLUE,       WHITE,     "Client Layer",
         "iOS App  |  Android App  |  Web PWA (React)"),
        (NAVY,       WHITE,     "API Gateway",
         "Auth Service  |  Closet Service  |  Outfit Agent  |  Recommendations  |  B2B Data API"),
        (RGBColor(0, 90, 160),  WHITE,  "AI & Data",
         "LLM with Vision (Outfit Agent)  |  Recommendation Engine  |  Data Pipeline"),
        (RGBColor(70, 130, 180), WHITE, "External Integrations",
         "Retailer OAuth APIs  |  Weather API  |  Google/Apple Calendar  |  Affiliate Networks"),
        (RGBColor(100, 100, 120), WHITE, "Data Store",
         "User & Closet DB  |  Product Catalog  |  B2B Data Product"),
    ]

    box_h = Inches(0.88)
    gap   = Inches(0.08)
    box_w = Inches(12.0)
    left  = Inches(0.65)
    start_top = Inches(1.3)

    for i, (fill, txt_color, label, content) in enumerate(layers):
        top = start_top + i * (box_h + gap)

        # Label strip
        lbl = slide.shapes.add_shape(1, left, top, Inches(1.9), box_h)
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = fill
        lbl.line.fill.background()
        tf = lbl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = txt_color

        # Content box
        content_left = left + Inches(2.0)
        content_w    = box_w - Inches(2.0)
        cb = slide.shapes.add_shape(1, content_left, top, content_w, box_h)
        cb.fill.solid()
        cb.fill.fore_color.rgb = RGBColor(230, 240, 250)
        cb.line.color.rgb = fill
        tf2 = cb.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = content
        run2.font.size = Pt(13)
        run2.font.color.rgb = NAVY


def slide_07_ai_agent(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "AI Outfit Agent — How It Thinks")
    add_divider(slide, Inches(1.05))

    flow_left  = Inches(1.5)
    flow_w     = Inches(8.5)
    box_h      = Inches(1.4)
    arrow_h    = Inches(0.4)

    tops = [Inches(1.3), Inches(3.1), Inches(4.9)]
    fills   = [LIGHT_GRAY, BLUE, LIGHT_GRAY]
    t_colors = [NAVY,      WHITE, NAVY]
    items = [
        ("Contextual Inputs",
         "Weather API  ·  Day of Week  ·  Occasion  ·  Calendar Events  ·  User Style Profile"),
        ("LLM with Vision",
         "Scans closet JSON + images  →  Applies style rules + color theory  →  "
         "Composes outfit per slot  →  Generates color palette"),
        ("Structured Output",
         "3–6 outfit options  ·  Item IDs per slot  ·  Color palette hex array  ·  "
         "Palette name  ·  Occasion tag  ·  Natural language styling note"),
    ]

    for i, ((label, body), top, fill, tc) in enumerate(zip(items, tops, fills, t_colors)):
        bg = slide.shapes.add_shape(1, flow_left, top, flow_w, box_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = BLUE

        add_textbox(slide, label,
                    flow_left + Inches(0.15), top + Inches(0.05),
                    flow_w - Inches(0.3), Inches(0.4),
                    font_size=14, bold=True, color=tc)

        add_textbox(slide, body,
                    flow_left + Inches(0.15), top + Inches(0.5),
                    flow_w - Inches(0.3), Inches(0.8),
                    font_size=12, color=tc, wrap=True)

        # Arrow
        if i < 2:
            add_textbox(slide, "↓",
                        flow_left + flow_w / 2 - Inches(0.2),
                        top + box_h,
                        Inches(0.4), arrow_h,
                        font_size=20, bold=True, color=BLUE,
                        align=PP_ALIGN.CENTER)

    # Side note
    side_left = Inches(10.3)
    side_top  = Inches(1.3)
    side_w    = Inches(2.7)
    side_h    = Inches(5.0)
    note_bg = slide.shapes.add_shape(1, side_left, side_top, side_w, side_h)
    note_bg.fill.solid()
    note_bg.fill.fore_color.rgb = LIGHT_GRAY
    note_bg.line.color.rgb = BLUE

    note_text = (
        "Session Rules\n\n"
        "• Min 3, Max 6 options per session\n\n"
        "• Each anchored to distinct color palette\n\n"
        "• Regeneration without losing prior options"
    )
    add_textbox(slide, note_text,
                side_left + Inches(0.15), side_top + Inches(0.2),
                side_w - Inches(0.3), side_h - Inches(0.4),
                font_size=12, color=NAVY, wrap=True)


def slide_08_business_model(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "Business Model — Three Revenue Streams")
    add_divider(slide, Inches(1.05))

    streams = [
        {
            "header": "Affiliate Commerce",
            "body": (
                "Commission on product purchases made through in-app suggestion links. "
                "Deep links to retailer pages with tracking parameters."
            ),
            "target": "Target: >8% click-through rate",
        },
        {
            "header": "B2B Data Licensing",
            "body": (
                "Structured consumer behavior data product licensed to fashion brands "
                "via AI Product Designer platform API."
            ),
            "target": "Target: 5 B2B clients Year 1",
        },
        {
            "header": "Future: Premium Subscription",
            "body": (
                "Advanced features: calendar integration, style coaching, trend alerts. "
                "Out of scope for v1."
            ),
            "target": "Note: V2 consideration",
        },
    ]

    box_w = Inches(3.8)
    box_h = Inches(4.5)
    gap   = Inches(0.3)
    start_left = Inches(0.55)
    top   = Inches(1.5)

    for i, s in enumerate(streams):
        left = start_left + i * (box_w + gap)

        # Box
        bg = slide.shapes.add_shape(1, left, top, box_w, box_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = BLUE

        # Header
        hdr = slide.shapes.add_shape(1, left, top, box_w, Inches(0.55))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = BLUE
        hdr.line.fill.background()
        tf = hdr.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = s["header"]
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Body
        add_textbox(slide, s["body"],
                    left + Inches(0.15), top + Inches(0.7),
                    box_w - Inches(0.3), Inches(3.0),
                    font_size=14, color=BODY_TEXT, wrap=True)

        # Target chip
        chip = slide.shapes.add_shape(1,
                left + Inches(0.15),
                top + box_h - Inches(0.55),
                box_w - Inches(0.3), Inches(0.4))
        chip.fill.solid()
        chip.fill.fore_color.rgb = BLUE
        chip.line.fill.background()
        tf2 = chip.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = s["target"]
        run2.font.size = Pt(12)
        run2.font.bold = True
        run2.font.color.rgb = WHITE


def slide_09_data_product(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "The Data Asset — AI Product Designer Pipeline")
    add_divider(slide, Inches(1.05))

    # Left column
    left_l = Inches(0.5)
    left_w = Inches(6.2)

    add_textbox(slide,
                "Every user interaction builds a rich, consented data product. "
                "Fashion brands query real consumer behavior — not surveys.",
                left_l, Inches(1.4), left_w, Inches(0.9),
                font_size=14, color=BODY_TEXT, wrap=True)

    add_textbox(slide, "Data Collected",
                left_l, Inches(2.45), left_w, Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    collected = (
        "• Purchase history & wardrobe composition\n"
        "• Outfit acceptance & color palette preferences\n"
        "• Product click & affiliate purchase behavior\n"
        "• AI-inferred lifestyle attributes & style archetype"
    )
    add_textbox(slide, collected,
                left_l, Inches(2.9), left_w, Inches(1.6),
                font_size=13, color=BODY_TEXT, wrap=True)

    add_textbox(slide, "B2B Schema Highlights",
                left_l, Inches(4.6), left_w, Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    schema = (
        "style_archetype  ·  price_tier  ·  color_palette_preference  ·  "
        "category_affinity  ·  product_taxonomy  ·  geo_region"
    )
    add_textbox(slide, schema,
                left_l, Inches(5.05), left_w, Inches(0.6),
                font_size=13, color=BODY_TEXT, italic=True, wrap=True)

    # Privacy statement
    priv_bg = slide.shapes.add_shape(1,
        left_l, Inches(5.8), left_w, Inches(0.9))
    priv_bg.fill.solid()
    priv_bg.fill.fore_color.rgb = RGBColor(230, 245, 230)
    priv_bg.line.color.rgb = GREEN
    add_textbox(slide,
                "GDPR & CCPA compliant  ·  User opt-in required  ·  "
                "Individual data withdrawal within 72 hours  ·  Privacy dashboard in-app",
                left_l + Inches(0.1), Inches(5.85), left_w - Inches(0.2), Inches(0.8),
                font_size=12, color=GREEN, bold=True, wrap=True)

    # Right column — delivery options
    right_l = Inches(7.2)
    right_w = Inches(5.7)

    add_textbox(slide, "Delivery Options",
                right_l, Inches(1.4), right_w, Inches(0.4),
                font_size=14, bold=True, color=NAVY)

    delivery_items = [
        "REST API (OAuth 2.0)",
        "Daily batch export (JSONL / CSV)",
        "Nightly profile refresh",
        "99.5% API uptime SLA",
    ]
    for j, item in enumerate(delivery_items):
        top_d = Inches(1.95) + j * Inches(0.75)
        chip = slide.shapes.add_shape(1, right_l, top_d, right_w, Inches(0.6))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(230, 240, 250)
        chip.line.color.rgb = BLUE
        tf = chip.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "  " + item
        run.font.size = Pt(14)
        run.font.color.rgb = NAVY


def slide_10_gtm(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "Go-to-Market Strategy")
    add_divider(slide, Inches(1.05))

    phases = [
        {
            "phase": "Phase 1",
            "label": "Launch",
            "body": (
                "Target style-conscious professionals 25–42 in top 5 US metros. "
                "App Store + Google Play launch. Seed with 10+ retailer API integrations "
                "covering majority of target user shopping behavior."
            ),
        },
        {
            "phase": "Phase 2",
            "label": "Grow",
            "body": (
                "Influencer partnerships in fashion/lifestyle vertical. "
                "Content marketing: 'style with what you own' positioning. "
                "Referral mechanics in-app. Target: 50K MAU Year 1."
            ),
        },
        {
            "phase": "Phase 3",
            "label": "Monetize & Scale",
            "body": (
                "Activate B2B data product at 10K+ user threshold. "
                "Launch affiliate network integrations. "
                "Expand retailer coverage internationally. "
                "Pipeline to AI Product Designer platform."
            ),
        },
    ]

    box_w = Inches(3.8)
    box_h = Inches(3.8)
    gap   = Inches(0.3)
    start_left = Inches(0.55)
    top   = Inches(1.5)

    for i, ph in enumerate(phases):
        left = start_left + i * (box_w + gap)

        bg = slide.shapes.add_shape(1, left, top, box_w, box_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = BLUE

        # Phase number pill
        pill = slide.shapes.add_shape(1, left + Inches(0.15), top + Inches(0.15),
                                      Inches(1.0), Inches(0.38))
        pill.fill.solid()
        pill.fill.fore_color.rgb = BLUE
        pill.line.fill.background()
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ph["phase"]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Label
        add_textbox(slide, ph["label"],
                    left + Inches(0.15), top + Inches(0.65),
                    box_w - Inches(0.3), Inches(0.5),
                    font_size=16, bold=True, color=NAVY)

        # Body
        add_textbox(slide, ph["body"],
                    left + Inches(0.15), top + Inches(1.25),
                    box_w - Inches(0.3), Inches(2.4),
                    font_size=13, color=BODY_TEXT, wrap=True)

    # Partnership row
    add_textbox(slide, "Key Partnership Types:",
                Inches(0.55), Inches(5.6), Inches(3.0), Inches(0.4),
                font_size=13, bold=True, color=NAVY)

    partners = [
        "Retailer API Partnerships",
        "Affiliate Networks",
        "Fashion Media / Influencers",
    ]
    for k, pt in enumerate(partners):
        chip = slide.shapes.add_shape(1,
            Inches(3.7) + k * Inches(2.9), Inches(5.55),
            Inches(2.7), Inches(0.45))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(210, 234, 255)
        chip.line.color.rgb = BLUE
        tf = chip.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = pt
        run.font.size = Pt(12)
        run.font.color.rgb = NAVY


def slide_11_metrics(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "Success Metrics — Year 1 Targets")
    add_divider(slide, Inches(1.05))

    rows = [
        ("Monthly Active Users",         "50,000",          "Validates consumer product-market fit"),
        ("Wardrobe Items Digitized",      ">30 per user",    "Signals engagement depth and data quality"),
        ("Outfit Sessions / Week",        ">3 per active user", "Core engagement loop health"),
        ("Affiliate CTR",                 ">8%",             "Revenue stream viability"),
        ("B2B Data Clients",              "5 clients",       "B2B revenue and strategic positioning"),
        ("User Consent Rate",             ">60%",            "Required for B2B product to function"),
    ]

    col_widths = [Inches(4.0), Inches(2.5), Inches(5.5)]
    col_lefts  = [Inches(0.5), Inches(4.6), Inches(7.2)]
    row_h      = Inches(0.62)
    hdr_top    = Inches(1.3)
    headers    = ["Metric", "Target", "Why It Matters"]

    # Header row
    for col_i, (hdr, cw, cl) in enumerate(zip(headers, col_widths, col_lefts)):
        hdr_bg = slide.shapes.add_shape(1, cl, hdr_top, cw, row_h)
        hdr_bg.fill.solid()
        hdr_bg.fill.fore_color.rgb = BLUE
        hdr_bg.line.fill.background()
        tf = hdr_bg.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Data rows
    for row_i, row in enumerate(rows):
        top = hdr_top + (row_i + 1) * row_h
        fill = WHITE if row_i % 2 == 0 else LIGHT_GRAY
        for col_i, (val, cw, cl) in enumerate(zip(row, col_widths, col_lefts)):
            cell_bg = slide.shapes.add_shape(1, cl, top, cw, row_h)
            cell_bg.fill.solid()
            cell_bg.fill.fore_color.rgb = fill
            cell_bg.line.color.rgb = RGBColor(200, 200, 200)

            align = PP_ALIGN.CENTER if col_i == 1 else PP_ALIGN.LEFT
            fc = GREEN if col_i == 1 else BODY_TEXT
            offset = Inches(0.1) if col_i != 1 else Inches(0)
            add_textbox(slide, val,
                        cl + offset, top + Inches(0.08),
                        cw - offset, row_h - Inches(0.1),
                        font_size=13,
                        bold=(col_i == 1),
                        color=fc,
                        align=align)


def slide_12_roadmap(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "Product Roadmap")
    add_divider(slide, Inches(1.05))

    columns = [
        {
            "label": "V1 — Now",
            "fill": BLUE,
            "items": [
                "Email + Google + Apple auth",
                "Retailer API wardrobe import",
                "Manual item upload with AI categorization",
                "AI outfit generation (3–6 options)",
                "Weather + occasion context",
                "Affiliate product suggestions",
                "B2B data pipeline (anonymized)",
            ],
        },
        {
            "label": "V2 — Next",
            "fill": NAVY,
            "items": [
                "Google/Apple Calendar integration",
                "Individual data consent + Privacy Dashboard",
                "Subscription tier with advanced features",
                "Menswear / gender-neutral taxonomy",
                "Social sharing of outfits",
                "In-app direct checkout",
            ],
        },
        {
            "label": "V3 — Future",
            "fill": RGBColor(70, 130, 180),
            "items": [
                "AI Product Designer platform (separate product)",
                "Style coaching / human stylist integration",
                "International retailer expansion",
                "Trend forecasting for B2B clients",
            ],
        },
    ]

    col_w = Inches(3.8)
    gap   = Inches(0.3)
    col_h = Inches(5.4)
    start_left = Inches(0.55)
    top   = Inches(1.3)

    for i, col in enumerate(columns):
        left = start_left + i * (col_w + gap)

        # Column bg
        bg = slide.shapes.add_shape(1, left, top, col_w, col_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = col["fill"]

        # Header
        hdr = slide.shapes.add_shape(1, left, top, col_w, Inches(0.55))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = col["fill"]
        hdr.line.fill.background()
        tf = hdr.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = col["label"]
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Items
        for j, item in enumerate(col["items"]):
            item_top = top + Inches(0.65) + j * Inches(0.68)
            item_bg = slide.shapes.add_shape(1,
                left + Inches(0.12), item_top,
                col_w - Inches(0.24), Inches(0.6))
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = WHITE
            item_bg.line.color.rgb = RGBColor(210, 220, 235)
            add_textbox(slide, "✓  " + item,
                        left + Inches(0.15), item_top + Inches(0.03),
                        col_w - Inches(0.3), Inches(0.56),
                        font_size=12, color=BODY_TEXT, wrap=True)


def slide_13_open_questions(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    add_bottom_border(slide)
    slide_title(slide, "Open Questions & Decisions Pending")
    add_divider(slide, Inches(1.05))

    questions = [
        ("1", "Which retailer API partnerships are confirmed at launch vs. requiring fallback credential flow?",
         "Business Development"),
        ("2", "What is the minimum viable brand list for launch to cover majority of target user shopping behavior?",
         "Product"),
        ("3", "Calendar integration (Google / Apple): in scope for v1 or deferred to v2?",
         "Product"),
        ("4", "Which affiliate network(s) will be used for product recommendation monetization?",
         "Business Development"),
        ("5", "What is the pricing model for the B2B data product? Per seat, per API call, or per cohort?",
         "Business Development"),
        ("6", "GDPR data processing agreements required for each B2B client — is legal review resourced?",
         "Legal"),
        ("7", "Computer vision model for manual item upload: build, buy, or use existing API (e.g., Google Vision)?",
         "Engineering"),
    ]

    col_widths = [Inches(0.5), Inches(8.5), Inches(2.8)]
    col_lefts  = [Inches(0.5), Inches(1.1), Inches(9.7)]
    row_h      = Inches(0.56)
    hdr_top    = Inches(1.25)
    headers    = ["#", "Question", "Owner"]

    # Header
    for hdr, cw, cl in zip(headers, col_widths, col_lefts):
        hdr_bg = slide.shapes.add_shape(1, cl, hdr_top, cw, row_h)
        hdr_bg.fill.solid()
        hdr_bg.fill.fore_color.rgb = BLUE
        hdr_bg.line.fill.background()
        tf = hdr_bg.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

    owner_colors = {
        "Business Development": RGBColor(0, 80, 160),
        "Product":              BLUE,
        "Legal":                RGBColor(120, 0, 120),
        "Engineering":          RGBColor(0, 100, 80),
    }

    for row_i, (num, question, owner) in enumerate(questions):
        top  = hdr_top + (row_i + 1) * row_h
        fill = WHITE if row_i % 2 == 0 else LIGHT_GRAY
        row_data = [num, question, owner]

        for col_i, (val, cw, cl) in enumerate(zip(row_data, col_widths, col_lefts)):
            cell_bg = slide.shapes.add_shape(1, cl, top, cw, row_h)
            cell_bg.fill.solid()
            cell_bg.fill.fore_color.rgb = fill
            cell_bg.line.color.rgb = RGBColor(200, 200, 200)

            if col_i == 0:
                # Number
                add_textbox(slide, val, cl, top + Inches(0.05),
                            cw, row_h - Inches(0.1),
                            font_size=13, bold=True, color=BLUE,
                            align=PP_ALIGN.CENTER)
            elif col_i == 1:
                add_textbox(slide, val, cl + Inches(0.08), top + Inches(0.05),
                            cw - Inches(0.1), row_h - Inches(0.1),
                            font_size=12, color=BODY_TEXT, wrap=True)
            else:
                oc = owner_colors.get(val, NAVY)
                chip = slide.shapes.add_shape(1, cl + Inches(0.05), top + Inches(0.1),
                                              cw - Inches(0.1), row_h - Inches(0.2))
                chip.fill.solid()
                chip.fill.fore_color.rgb = RGBColor(230, 240, 250)
                chip.line.color.rgb = oc
                tf = chip.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = val
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = oc


def slide_14_closing(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BLUE)

    # Decorative circle (large, subtle)
    circ = slide.shapes.add_shape(9,
        Inches(8.5), Inches(-1.5), Inches(6), Inches(6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0, 100, 180)
    circ.line.fill.background()

    add_textbox(slide,
                "AI Stylist From Your Closet",
                Inches(1), Inches(2.0), Inches(11.33), Inches(1.4),
                font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Built on the belief that the best outfit is already in your closet.",
                Inches(1), Inches(3.7), Inches(11.33), Inches(0.7),
                font_size=20, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # Divider
    bar = slide.shapes.add_shape(1, Inches(4.0), Inches(4.7), Inches(5.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()

    add_textbox(slide,
                "Julia Lee  ·  juliahwanlee@gmail.com  ·  linkedin.com/in/juliahwanlee",
                Inches(1), Inches(4.9), Inches(11.33), Inches(0.5),
                font_size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    print("Building slide 1  — Title")
    slide_01_title(prs)
    print("Building slide 2  — The Problem")
    slide_02_problem(prs)
    print("Building slide 3  — The Solution")
    slide_03_solution(prs)
    print("Building slide 4  — How It Works")
    slide_04_how_it_works(prs)
    print("Building slide 5  — User Personas")
    slide_05_personas(prs)
    print("Building slide 6  — Product Architecture")
    slide_06_architecture(prs)
    print("Building slide 7  — AI Agent Architecture")
    slide_07_ai_agent(prs)
    print("Building slide 8  — Business Model")
    slide_08_business_model(prs)
    print("Building slide 9  — B2B Data Product")
    slide_09_data_product(prs)
    print("Building slide 10 — Go-to-Market Strategy")
    slide_10_gtm(prs)
    print("Building slide 11 — Success Metrics")
    slide_11_metrics(prs)
    print("Building slide 12 — Product Roadmap")
    slide_12_roadmap(prs)
    print("Building slide 13 — Open Questions")
    slide_13_open_questions(prs)
    print("Building slide 14 — Closing")
    slide_14_closing(prs)

    output_path = "/Users/julialee/ClosetFinder/closetfinder_deck.pptx"
    prs.save(output_path)
    print(f"\nDeck saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
